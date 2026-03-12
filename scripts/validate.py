#!/usr/bin/env python3
"""
PPTX validation script.

Validates generated PPTX files using:
1. ZIP validity check
2. OPC structure check ([Content_Types].xml, .rels files)
3. python-pptx round-trip (structural validation)
4. LibreOffice headless PDF conversion (rendering validation)

Usage:
    python3 scripts/validate.py <pptx_path> [--slides N] [--libreoffice]

Exit codes:
    0: All checks passed
    1: Validation failed
"""

import json
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

LIBREOFFICE_PATH = Path(__file__).with_name("libreoffice")


def validate_zip(path: str) -> list[str]:
    """Verify the file is a valid ZIP."""
    errors = []
    try:
        with zipfile.ZipFile(path, "r") as zf:
            bad = zf.testzip()
            if bad is not None:
                errors.append(f"Corrupt ZIP entry: {bad}")
    except zipfile.BadZipFile as e:
        errors.append(f"Invalid ZIP file: {e}")
    return errors


def validate_opc_structure(path: str) -> list[str]:
    """Verify OPC required files are present and well-formed XML."""
    errors = []
    required_files = ["[Content_Types].xml"]

    with zipfile.ZipFile(path, "r") as zf:
        names = zf.namelist()
        for req in required_files:
            if req not in names:
                errors.append(f"Missing required file: {req}")

        # Check all XML files are well-formed
        for name in names:
            if name.endswith(".xml") or name.endswith(".rels"):
                try:
                    data = zf.read(name)
                    etree.fromstring(data)
                except etree.XMLSyntaxError as e:
                    errors.append(f"Malformed XML in {name}: {e}")

        # Check for .rels files
        if "_rels/.rels" not in names:
            errors.append("Missing _rels/.rels")

    return errors


def validate_python_pptx(path: str, expected_slides: int | None = None) -> list[str]:
    """Round-trip through python-pptx and verify content."""
    errors = []
    try:
        prs = Presentation(path)
        slide_count = len(prs.slides)
        if expected_slides is not None and slide_count != expected_slides:
            errors.append(
                f"Expected {expected_slides} slides, got {slide_count}"
            )

        # Verify each slide is accessible
        for i, slide in enumerate(prs.slides):
            try:
                _ = len(slide.shapes)
            except Exception as e:
                errors.append(f"Error reading slide {i + 1}: {e}")

        # Output structured info
        info = {
            "slide_count": slide_count,
            "slides": [],
        }
        for i, slide in enumerate(prs.slides):
            slide_info = {
                "index": i,
                "shape_count": len(slide.shapes),
                "shapes": [],
            }
            for shape in slide.shapes:
                shape_info = {
                    "name": shape.name,
                    "shape_type": str(shape.shape_type) if shape.shape_type else None,
                    "has_text_frame": shape.has_text_frame,
                }
                if shape.has_text_frame:
                    shape_info["text"] = shape.text_frame.text

                # Image detection
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_info["is_picture"] = True
                    try:
                        shape_info["image_content_type"] = shape.image.content_type
                    except Exception:
                        pass

                # Table detection
                if shape.has_table:
                    tbl = shape.table
                    shape_info["is_table"] = True
                    shape_info["table_rows"] = len(tbl.rows)
                    shape_info["table_cols"] = len(tbl.columns)
                    table_data = []
                    for row in tbl.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        table_data.append(row_data)
                    shape_info["table_data"] = table_data

                # Chart detection
                if getattr(shape, "has_chart", False):
                    shape_info["is_chart"] = True

                slide_info["shapes"].append(shape_info)
            info["slides"].append(slide_info)

        # Print info as JSON for test consumption
        print(json.dumps(info))

    except Exception as e:
        errors.append(f"python-pptx failed to open: {e}")

    return errors


def validate_libreoffice(path: str) -> list[str]:
    """Convert to PDF with LibreOffice headless."""
    errors = []
    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            result = subprocess.run(
                [
                    str(LIBREOFFICE_PATH),
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmpdir,
                    path,
                ],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode != 0:
                errors.append(
                    f"LibreOffice conversion failed: {result.stderr}"
                )
            else:
                # Check PDF was actually created
                pdf_files = list(Path(tmpdir).glob("*.pdf"))
                if not pdf_files:
                    errors.append("LibreOffice produced no PDF output")
        except subprocess.TimeoutExpired:
            errors.append("LibreOffice conversion timed out")
        except FileNotFoundError:
            errors.append("LibreOffice not found")

    return errors


def main() -> int:
    import argparse

    parser = argparse.ArgumentParser(description="Validate PPTX files")
    parser.add_argument("pptx_path", help="Path to PPTX file")
    parser.add_argument(
        "--slides", type=int, default=None, help="Expected slide count"
    )
    parser.add_argument(
        "--libreoffice", action="store_true", help="Run LibreOffice validation"
    )
    args = parser.parse_args()

    all_errors: list[str] = []

    # ZIP check
    errors = validate_zip(args.pptx_path)
    all_errors.extend(errors)
    if errors:
        print(f"ZIP validation FAILED: {errors}", file=sys.stderr)
        return 1

    # OPC structure check
    errors = validate_opc_structure(args.pptx_path)
    all_errors.extend(errors)
    if errors:
        print(f"OPC validation FAILED: {errors}", file=sys.stderr)

    # python-pptx round-trip
    errors = validate_python_pptx(args.pptx_path, args.slides)
    all_errors.extend(errors)
    if errors:
        print(f"python-pptx validation FAILED: {errors}", file=sys.stderr)

    # LibreOffice (optional)
    if args.libreoffice:
        errors = validate_libreoffice(args.pptx_path)
        all_errors.extend(errors)
        if errors:
            print(f"LibreOffice validation FAILED: {errors}", file=sys.stderr)

    if all_errors:
        return 1

    return 0


if __name__ == "__main__":
    sys.exit(main())
